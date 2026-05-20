"""슬라이드 자료(PPTX) 산출물.

lecture-postpro-mcp 의 make_cardnews() 가 그대로 import해서 쓸 수 있도록
`make_slides_lite()` 를 Streamlit/RAG 비의존 **순수 async 함수**로 분리해 둔다.
PPTX 렌더링(`_render_pptx`)도 RAG/Streamlit과 독립이라 MCP 측에서 그대로 호출 가능.
"""
from __future__ import annotations

import json
import re
from datetime import datetime
from pathlib import Path
from typing import Awaitable, Callable, Literal

from core.llm_client import get_llm_func
from core.ingest.subtitle import parse_subtitle_text

from ._base import ArtifactMeta, ArtifactResult, load_prompt

META = ArtifactMeta(
    key="slides",
    title="PPTX 보고서",
    icon="📊",
    order=20,
    model_profile="creative",
    accepts=("subtitle", "text"),
    description="자막/텍스트 → 한국어 발표용 PPTX (6장 디폴트). 회사 양식 자동 적용.",
)


LLMFunc = Callable[..., Awaitable[str]]

# ── 고정 디자인 테마 ─────────────────────────────────────────────────
# 비즈니스 친화 컬러. 필요 시 한 곳만 바꾸면 전체 슬라이드에 반영.
THEME = {
    "bg_main": (0xF7, 0xF8, 0xFB),       # 본문 배경 - 매우 옅은 회색
    "bg_title": (0x1F, 0x2A, 0x44),      # 표지 배경 - 짙은 네이비
    "accent": (0x2E, 0x5B, 0xFF),        # 강조 - 블루
    "text_dark": (0x18, 0x1B, 0x2C),     # 본문 글자
    "text_light": (0xFF, 0xFF, 0xFF),    # 다크 배경 위 글자
    "text_muted": (0x6B, 0x72, 0x80),    # 보조 글자
    "bullet": (0x2E, 0x5B, 0xFF),        # 불릿 색
    "font_title": "맑은 고딕",
    "font_body": "맑은 고딕",
}


def _strip_fence(text: str) -> str:
    m = re.search(r"```(?:json)?\s*([\s\S]+?)```", text)
    return (m.group(1) if m else text).strip()


def _to_markdown_preview(bundle: dict) -> str:
    """Streamlit 미리보기용 마크다운 (PPTX 다운로드 전 텍스트 확인용)."""
    lines: list[str] = []
    title = bundle.get("title", "")
    subtitle = bundle.get("subtitle", "")
    if title:
        lines.append(f"# {title}")
    if subtitle:
        lines.append(f"_{subtitle}_")
        lines.append("")

    for i, slide in enumerate(bundle.get("slides", []), 1):
        lines.append(f"## {i}. {slide.get('title','')}")
        for b in slide.get("bullets", []):
            lines.append(f"- {b}")
        notes = slide.get("speaker_notes")
        if notes:
            lines.append("")
            lines.append(f"> 🎤 _{notes}_")
        kws = slide.get("keywords") or []
        if kws:
            lines.append("")
            lines.append("**키워드**: " + ", ".join(kws))
        lines.append("")
    return "\n".join(lines).strip()


async def make_slides_lite(
    subtitle_text: str,
    subtitle_format: Literal["srt", "vtt", "plain"] = "plain",
    metadata: dict | None = None,
    llm_call: LLMFunc | None = None,
    slides: int = 6,
    prompt_key: str = "slides",
) -> dict:
    """Pure async function — reused by lecture-postpro-mcp + slide_outline.

    prompt_key: prompts/<key>_ko.md 를 읽는다. PPTX 보고서는 "slides",
    슬라이드 교안(텍스트) 는 "slide_outline" 로 분리.

    Returns: {"title", "subtitle", "slides":[{title, bullets, speaker_notes, keywords}], "markdown": "..."}
    """
    clean = (
        parse_subtitle_text(subtitle_text, subtitle_format)
        if subtitle_format in ("srt", "vtt")
        else subtitle_text
    )

    prompt_tpl = load_prompt(prompt_key) or (
        "다음 강의 자막을 바탕으로 {{N}}개의 슬라이드로 정리하라."
    )
    prompt = (
        prompt_tpl.replace("{{N}}", str(slides))
        .replace("{{N-1}}", str(slides - 1))
        + "\n\n[메타데이터]\n"
        + json.dumps(metadata or {}, ensure_ascii=False)
        + "\n\n[자막]\n"
        + clean
    )

    caller = llm_call or get_llm_func("creative")
    raw = await caller(prompt)
    payload = _strip_fence(raw)
    try:
        bundle = json.loads(payload)
    except json.JSONDecodeError:
        bundle = {
            "title": "파싱 실패",
            "subtitle": "",
            "slides": [{"title": "원본", "bullets": [raw[:200]], "speaker_notes": "", "keywords": []}],
        }

    bundle.setdefault("title", "발표 자료")
    bundle.setdefault("subtitle", "")
    bundle.setdefault("slides", [])
    bundle["markdown"] = _to_markdown_preview(bundle)
    return bundle


# ── PPTX 렌더링 ───────────────────────────────────────────────────────
def _resolve_template_path() -> Path | None:
    """회사 양식 파일 경로. 우선순위:
    1) 환경변수 PPTX_TEMPLATE_PATH (절대/상대 경로)
    2) 프로젝트 루트의 assets/pptx_template.pptx
    파일이 실제 존재할 때만 반환, 없으면 None → 빈 Presentation 으로 폴백.
    스펙: docs/pptx_template_spec.md
    """
    import os
    env_path = os.environ.get("PPTX_TEMPLATE_PATH", "").strip()
    if env_path:
        p = Path(env_path)
        if p.exists():
            return p
    default = Path(__file__).resolve().parent.parent / "assets" / "pptx_template.pptx"
    return default if default.exists() else None


def _render_pptx(bundle: dict, out_path: Path) -> Path:
    """python-pptx 로 PPTX 파일 생성. 회사 양식(있으면) 우선, 없으면 내장 테마."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    out_path.parent.mkdir(parents=True, exist_ok=True)

    template_path = _resolve_template_path()
    if template_path:
        prs = Presentation(str(template_path))
        # 양식에 샘플 슬라이드가 들어있을 수 있어 모두 비운다 — 마스터/레이아웃만 유지.
        xml_slides = prs.slides._sldIdLst  # noqa: SLF001
        for sld in list(xml_slides):
            xml_slides.remove(sld)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)   # 16:9
        prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height

    # 양식이 있으면 마스터의 첫 두 레이아웃(표지/본문)을 사용. 없으면 blank=레이아웃 6.
    blank = prs.slide_layouts[6] if not template_path else prs.slide_layouts[
        min(1, len(prs.slide_layouts) - 1)
    ]

    # ─ 표지 ─
    title_slide = prs.slides.add_slide(blank)
    _fill_bg(title_slide, RGBColor(*THEME["bg_title"]))
    _add_text(
        title_slide,
        bundle.get("title", "발표 자료"),
        left=Inches(0.7), top=Inches(2.6), width=SW - Inches(1.4), height=Inches(1.5),
        size=44, bold=True, color=THEME["text_light"], font=THEME["font_title"],
    )
    if bundle.get("subtitle"):
        _add_text(
            title_slide,
            bundle["subtitle"],
            left=Inches(0.7), top=Inches(4.2), width=SW - Inches(1.4), height=Inches(0.7),
            size=20, bold=False, color=THEME["text_muted"], font=THEME["font_body"],
        )
    _add_bar(title_slide, left=Inches(0.7), top=Inches(4.05), width=Inches(1.2), height=Emu(45720),
             color=THEME["accent"])

    # ─ 본문 슬라이드들 ─
    total = len(bundle.get("slides", []))
    for idx, slide_data in enumerate(bundle.get("slides", []), 1):
        slide = prs.slides.add_slide(blank)
        _fill_bg(slide, RGBColor(*THEME["bg_main"]))

        # 상단 강조 바
        _add_bar(slide, left=Inches(0.5), top=Inches(0.5), width=Inches(0.6), height=Emu(57150),
                 color=THEME["accent"])

        # 타이틀
        _add_text(
            slide,
            slide_data.get("title", ""),
            left=Inches(1.25), top=Inches(0.35), width=SW - Inches(2), height=Inches(0.9),
            size=32, bold=True, color=THEME["text_dark"], font=THEME["font_title"],
        )

        # 불릿
        bullets = slide_data.get("bullets") or []
        if bullets:
            _add_bullets(
                slide,
                bullets,
                left=Inches(1.0), top=Inches(1.6), width=SW - Inches(2), height=Inches(4.8),
                color=THEME["text_dark"], bullet_color=THEME["bullet"], font=THEME["font_body"],
            )

        # 키워드 (하단)
        kws = slide_data.get("keywords") or []
        if kws:
            _add_text(
                slide,
                "  ·  ".join(f"#{k}" for k in kws),
                left=Inches(1.0), top=SH - Inches(0.95), width=SW - Inches(2), height=Inches(0.5),
                size=12, bold=False, color=THEME["text_muted"], font=THEME["font_body"],
            )

        # 페이지 번호
        _add_text(
            slide,
            f"{idx} / {total}",
            left=SW - Inches(1.3), top=SH - Inches(0.55), width=Inches(0.8), height=Inches(0.3),
            size=10, bold=False, color=THEME["text_muted"], font=THEME["font_body"],
            align_right=True,
        )

        # 스피커 노트
        notes = slide_data.get("speaker_notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    prs.save(out_path)
    return out_path


def _fill_bg(slide, rgb) -> None:
    from pptx.dml.color import RGBColor
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb if isinstance(rgb, RGBColor) else RGBColor(*rgb)


def _add_text(slide, text, *, left, top, width, height, size, bold, color, font, align_right=False):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if align_right:
        p.alignment = PP_ALIGN.RIGHT
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = RGBColor(*color)


def _add_bullets(slide, items, *, left, top, width, height, color, bullet_color, font):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ""
        # 컬러 불릿 마커 + 본문 두 개 run 으로 색 분리
        run_marker = p.add_run()
        run_marker.text = "●  "
        run_marker.font.name = font
        run_marker.font.size = Pt(20)
        run_marker.font.color.rgb = RGBColor(*bullet_color)

        run_body = p.add_run()
        run_body.text = item
        run_body.font.name = font
        run_body.font.size = Pt(20)
        run_body.font.color.rgb = RGBColor(*color)

        p.space_after = Pt(10)


def _add_bar(slide, *, left, top, width, height, color):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*color)
    shape.line.fill.background()


# ── Studio entrypoint ────────────────────────────────────────────────
async def generate(rag, context: dict) -> ArtifactResult:
    """Studio entrypoint.

    context: {'subtitle_text'?, 'subtitle_format'?, 'metadata'?, 'artifacts_dir'?, 'deck_name'?, 'slides'?}
    artifacts_dir 가 들어오면 그 안에 PPTX 저장 → result.files 에 path 포함 → UI 다운로드 버튼.
    """
    subtitle_text = context.get("subtitle_text") or ""
    if not subtitle_text:
        # 노트북 sources 에서 자막을 직접 읽어 합친다 (RAG 우회 — 카드뉴스와 동일 패턴).
        # 자막이 없는 노트북(예: PDF 전용)에서만 RAG 검색으로 폴백.
        notebook_name = context.get("notebook_name", "default")
        from core.rag import NotebookPaths
        from core.ingest.subtitle import parse_subtitle
        paths = NotebookPaths.for_notebook(notebook_name)
        subtitle_files = sorted([
            p for p in paths.sources.glob("*")
            if p.suffix.lower() in {".vtt", ".srt"}
        ])
        if subtitle_files:
            subtitle_text = "\n\n".join(
                f"[{p.stem}]\n{parse_subtitle(p)}" for p in subtitle_files
            )
        else:
            subtitle_text = await rag.aquery(
                "자막/대본 내용을 그대로 긴 문장으로 다시 출력하라. 요약 금지.",
                mode="hybrid",
                top_k=20,
            )

    bundle = await make_slides_lite(
        subtitle_text=subtitle_text,
        subtitle_format=context.get("subtitle_format", "plain"),
        metadata=context.get("metadata"),
        slides=int(context.get("slides", 6)),
    )

    files: list[Path] = []
    artifacts_dir = context.get("artifacts_dir")
    if artifacts_dir:
        out_dir = Path(artifacts_dir) / META.key
        deck_name = (context.get("deck_name") or bundle.get("title") or "slides").strip()
        safe = re.sub(r"[^\w\-가-힣]+", "_", deck_name)[:60] or "slides"
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        pptx_path = out_dir / f"{safe}_{stamp}.pptx"
        _render_pptx(bundle, pptx_path)
        files.append(pptx_path)

    return ArtifactResult(
        key=META.key,
        title=META.title,
        markdown=bundle.get("markdown", ""),
        data=bundle,
        files=files,
    )


def render(result: ArtifactResult) -> None:
    """Streamlit 표시 — 다운로드 버튼 + 미리보기 마크다운."""
    import streamlit as st

    if result.files:
        st.success(f"PPTX 생성 완료 — `{result.files[0].name}`")
    st.markdown(result.markdown or "_(빈 결과)_")
